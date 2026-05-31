import os
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Emu
import tempfile
BASE_DIR = Path(__file__).parent
PDF_PATH = BASE_DIR / "June2610X Luma - AI Visibility Deck.pdf"
PPTX_PATH = BASE_DIR / "June2610X Luma - AI Visibility Deck.pptx"
DPI = 200
def pdf_to_pptx(pdf_path, pptx_path, dpi=200):
    pages = convert_from_path(str(pdf_path), dpi=dpi)
    print(str(len(pages)) + " pages found.")
    w, h = pages[0].size
    sw = Emu(int(w * 914400 / dpi))
    sh = Emu(int(h * 914400 / dpi))
    prs = Presentation()
    prs.slide_width = sw
    prs.slide_height = sh
    blank = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as tmp:
        for i, pg in enumerate(pages, 1):
            print("  Slide " + str(i) + "/" + str(len(pages)))
            p = os.path.join(tmp, "p%04d.png" % i)
            pg.save(p, "PNG")
            sl = prs.slides.add_slide(blank)
            sl.shapes.add_picture(p, 0, 0, sw, sh)
    prs.save(str(pptx_path))
    print("Saved: " + str(pptx_path))
if __name__ == "__main__":
    pdf_to_pptx(PDF_PATH, PPTX_PATH, DPI)
